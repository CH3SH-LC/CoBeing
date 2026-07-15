from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # deep navy
BG_CARD = RGBColor(0x16, 0x21, 0x3A)       # card bg
ACCENT = RGBColor(0x00, 0xD4, 0xFF)        # cyan accent
ACCENT2 = RGBColor(0x7C, 0x5C, 0xFF)       # purple accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xB0, 0xC4, 0xDE)         # light steel blue
GRAY = RGBColor(0x6B, 0x7B, 0x93)
GREEN = RGBColor(0x00, 0xE6, 0x96)
ORANGE = RGBColor(0xFF, 0x9F, 0x43)
RED = RGBColor(0xFF, 0x6B, 0x6B)
HIGHLIGHT = RGBColor(0x00, 0xD4, 0xFF)

W = prs.slide_width
H = prs.slide_height


def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf


def add_para(tf, text, size=18, color=WHITE, bold=False, space_before=Pt(6),
             align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    if space_before:
        p.space_before = space_before
    return p


def add_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_line(slide, left, top, width, color=ACCENT, thickness=Pt(3)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_circle(slide, left, top, size, fill_color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


# ============================================================
# Slide 1 - Cover
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

# accent line
add_line(slide, Inches(1), Inches(2.8), Inches(2), ACCENT, Pt(4))

# title
add_text(slide, Inches(1), Inches(3.0), Inches(8), Inches(1.2),
         "CoBeing", size=56, color=WHITE, bold=True)
tf = add_text(slide, Inches(1), Inches(4.0), Inches(8), Inches(0.6),
              "多智能体协作框架", size=28, color=LIGHT)

# subtitle
add_text(slide, Inches(1), Inches(4.8), Inches(8), Inches(0.5),
        "让 AI Agents 组队干活", size=20, color=ACCENT)

# info
add_text(slide, Inches(1), Inches(5.8), Inches(6), Inches(0.4),
        '玄武区第二届"模数杯"AI创新挑战赛  |  AI智能体赛道', size=14, color=GRAY)
add_text(slide, Inches(1), Inches(6.2), Inches(6), Inches(0.4),
        "2026年6月  |  ch3sh_lc@sjtu.edu.cn", size=14, color=GRAY)

# decorative circle
add_circle(slide, Inches(10.5), Inches(1.5), Inches(3), RGBColor(0x0A, 0x2E, 0x50))
add_circle(slide, Inches(11.0), Inches(2.0), Inches(2), RGBColor(0x0D, 0x38, 0x60))

# ============================================================
# Slide 2 - Problem & Opportunity
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_line(slide, Inches(0.8), Inches(0.6), Inches(1.5), ACCENT)
add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.6),
         "问题与机会", size=36, color=WHITE, bold=True)

# quote box
add_rect(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.2), BG_CARD, ACCENT)
tf = add_text(slide, Inches(1.2), Inches(1.75), Inches(10.8), Inches(0.9),
              '"说一句帮我组建开发团队，AI自动创建需求分析师、架构师、工程师，\n'
              '像真实团队一样分工协作，交付成果。"',
              size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# 3 pain points
cards = [
    ("一问一答", "当前AI助手停留在单轮对话\n无法处理复杂任务"),
    ("预设流程", "现有框架靠串行执行\nAgent缺乏真正自主沟通"),
    ("零记忆", "没有持久化记忆\n每次会话从零开始"),
]
for i, (title, desc) in enumerate(cards):
    x = Inches(0.8 + i * 4.0)
    add_rect(slide, x, Inches(3.2), Inches(3.6), Inches(2.0), BG_CARD)
    add_circle(slide, x + Inches(0.3), Inches(3.45), Inches(0.5), ACCENT2)
    add_text(slide, x + Inches(1.0), Inches(3.4), Inches(2.4), Inches(0.4),
             title, size=20, color=WHITE, bold=True)
    add_text(slide, x + Inches(0.3), Inches(4.1), Inches(3.0), Inches(1.0),
             desc, size=15, color=LIGHT)

# data
add_rect(slide, Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.2), RGBColor(0x0A, 0x25, 0x40))
tf = add_text(slide, Inches(1.2), Inches(5.7), Inches(5), Inches(0.4),
              "市场数据", size=16, color=GREEN, bold=True)
add_para(tf, "2025年全球AI Agent市场 65亿美元  |  年增长 40%+  |  预计2028年达 260亿美元",
         size=15, color=LIGHT)

# ============================================================
# Slide 3 - Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_line(slide, Inches(0.8), Inches(0.6), Inches(1.5), ACCENT)
add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.6),
         "解决方案", size=36, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
         "CoBeing — 让AI Agent组队干活", size=22, color=ACCENT, bold=True)
add_text(slide, Inches(0.8), Inches(2.1), Inches(11), Inches(0.5),
         "多智能体协作框架，用户用自然语言描述需求，系统自动组建AI团队、分配任务、协作执行。",
         size=16, color=LIGHT)

# 4 cards
dims = [
    ("零学习成本", "对话即搭建\n管家帮你搞定一切", ACCENT),
    ("长效记忆", "Agent越用越懂你\n记忆跨对话不丢失", GREEN),
    ("自主协作", "Agent主动分工沟通\nTODO驱动自动交付", ORANGE),
    ("本地优先", "数据不上云\n隐私第一", ACCENT2),
]
for i, (title, desc, clr) in enumerate(dims):
    x = Inches(0.8 + i * 3.1)
    add_rect(slide, x, Inches(3.0), Inches(2.8), Inches(3.2), BG_CARD, clr)
    add_circle(slide, x + Inches(1.0), Inches(3.3), Inches(0.7), clr)
    add_text(slide, x + Inches(0.2), Inches(4.2), Inches(2.4), Inches(0.4),
             title, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(4.7), Inches(2.4), Inches(1.2),
             desc, size=14, color=LIGHT, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 4 - Demo
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_line(slide, Inches(0.8), Inches(0.6), Inches(1.5), ACCENT)
add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.6),
         "核心功能展示", size=36, color=WHITE, bold=True)

# flow steps
steps = [
    ("01", "用户发起请求", '"帮我组建一个开发团队\n做一个Todo应用"'),
    ("02", "管家分析需求", "推荐Agent角色\n创建团队成员"),
    ("03", "群主分解任务", "子任务分配\n依赖关系管理"),
    ("04", "Agent协作执行", "主动沟通 · 进度同步\n相互审核"),
    ("05", "成果交付", "代码 + 文档\n经验沉淀"),
]
for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.5 + i * 2.5)
    add_rect(slide, x, Inches(1.8), Inches(2.2), Inches(3.5), BG_CARD)
    # number circle
    add_circle(slide, x + Inches(0.7), Inches(2.0), Inches(0.7), ACCENT)
    add_text(slide, x + Inches(0.7), Inches(2.05), Inches(0.7), Inches(0.6),
             num, size=20, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(2.9), Inches(1.8), Inches(0.4),
             title, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(3.4), Inches(1.8), Inches(1.5),
             desc, size=13, color=LIGHT, align=PP_ALIGN.CENTER)

# arrow connectors
for i in range(4):
    x = Inches(2.7 + i * 2.5)
    add_text(slide, x, Inches(3.0), Inches(0.4), Inches(0.4),
             "→", size=24, color=ACCENT, bold=True)

# bottom note
add_text(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.4),
         "记忆不丢失，经验可复用 —— Agent越用越懂你",
         size=16, color=GREEN, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 5 - Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_line(slide, Inches(0.8), Inches(0.6), Inches(1.5), ACCENT)
add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.6),
         "技术架构", size=36, color=WHITE, bold=True)

# Architecture diagram using boxes
layers = [
    ("用户层", "GUI (React 19 + Tauri 2.0)  |  QQ Bot  |  MCP 插件", Inches(0.8), Inches(1.6), ACCENT),
    ("管家层", "Butler Agent — 自然语言理解 → 需求分析 → Agent/群组管理", Inches(0.8), Inches(2.7), GREEN),
    ("协作层", "群组工作组  |  WakeSystem事件队列  |  任务依赖链  |  投票表决", Inches(0.8), Inches(3.8), ORANGE),
    ("Agent层", "自治文件系统  |  对话循环  |  记忆引擎  |  技能系统", Inches(0.8), Inches(4.9), ACCENT2),
    ("基础设施", "LLM网关(7家厂商)  |  工具系统  |  Docker沙箱  |  MCP客户端", Inches(0.8), Inches(6.0), RGBColor(0xFF, 0x6B, 0x9D)),
]
for title, desc, x, y, clr in layers:
    add_rect(slide, x, y, Inches(8.5), Inches(0.9), BG_CARD, clr)
    add_text(slide, x + Inches(0.3), y + Inches(0.05), Inches(1.5), Inches(0.4),
             title, size=15, color=clr, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.45), Inches(7.8), Inches(0.4),
             desc, size=12, color=LIGHT)

# Stats on the right
add_rect(slide, Inches(9.8), Inches(1.6), Inches(3.0), Inches(5.3), BG_CARD)
stats = [
    ("~25K", "行TypeScript代码"),
    ("417", "个测试用例"),
    ("43", "个测试文件"),
    ("7", "家LLM厂商"),
    ("50+", "MCP工具"),
    ("5", "级权限系统"),
]
for i, (num, label) in enumerate(stats):
    y = Inches(1.8 + i * 0.85)
    add_text(slide, Inches(10.0), y, Inches(2.6), Inches(0.35),
             num, size=22, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(10.0), y + Inches(0.35), Inches(2.6), Inches(0.3),
             label, size=11, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 6 - Core Tech
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_line(slide, Inches(0.8), Inches(0.6), Inches(1.5), ACCENT)
add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.6),
         "三大技术支柱", size=36, color=WHITE, bold=True)

techs = [
    ("Agent自治系统", ACCENT, [
        "8文件自治文件系统",
        "SOUL / CHARACTER / JOB / EXPERIENCE",
        "System Prompt 分文件构建链",
        "Agent 通过对话自我进化",
    ]),
    ("事件驱动协作", GREEN, [
        "@mention 唤醒机制",
        "WakeSystem 事件队列",
        "三层上下文注入",
        "子任务依赖链 + 投票表决",
    ]),
    ("安全沙箱", ORANGE, [
        "Docker 容器隔离执行",
        "5级权限系统",
        "网络白名单 + 运行时检测",
        "API Key AES-256-GCM 加密",
    ]),
]
for i, (title, clr, items) in enumerate(techs):
    x = Inches(0.8 + i * 4.1)
    add_rect(slide, x, Inches(1.6), Inches(3.8), Inches(5.0), BG_CARD, clr)
    add_circle(slide, x + Inches(1.4), Inches(1.9), Inches(0.8), clr)
    add_text(slide, x + Inches(0.3), Inches(2.9), Inches(3.2), Inches(0.4),
             title, size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_line(slide, x + Inches(0.8), Inches(3.4), Inches(2.2), clr, Pt(2))
    for j, item in enumerate(items):
        add_text(slide, x + Inches(0.4), Inches(3.7 + j * 0.7), Inches(3.0), Inches(0.5),
                 f"•  {item}", size=13, color=LIGHT)

# ============================================================
# Slide 7 - Comparison
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_line(slide, Inches(0.8), Inches(0.6), Inches(1.5), ACCENT)
add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.6),
         "与现有方案有何不同？", size=36, color=WHITE, bold=True)

# Table header
headers = ["维度", "AutoGPT", "CrewAI", "MetaGPT", "CoBeing"]
col_w = [Inches(1.8), Inches(2.2), Inches(2.2), Inches(2.2), Inches(2.8)]
col_x = [Inches(0.8)]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)

y_start = Inches(1.6)
row_h = Inches(0.65)

# header row
for i, (hdr, x, w) in enumerate(zip(headers, col_x, col_w)):
    clr = ACCENT if i == 4 else RGBColor(0x1A, 0x2D, 0x4A)
    add_rect(slide, x, y_start, w, row_h, clr)
    add_text(slide, x, y_start + Inches(0.1), w, row_h,
             hdr, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

rows = [
    ["架构", "单Agent", "预设流程", "文档驱动", "事件驱动"],
    ["协作", "无", "串行执行", "角色固定", "自治群组"],
    ["记忆", "无持久化", "短期", "文档沉淀", "分层记忆+经验"],
    ["厂商", "绑定OpenAI", "绑定OpenAI", "绑定OpenAI", "7家国产"],
    ["GUI", "无", "无", "无", "完整桌面应用"],
]
for ri, row in enumerate(rows):
    y = y_start + row_h * (ri + 1)
    for ci, (cell, x, w) in enumerate(zip(row, col_x, col_w)):
        bg = BG_CARD if ci != 4 else RGBColor(0x0A, 0x2E, 0x50)
        add_rect(slide, x, y, w, row_h, bg)
        clr = ACCENT if ci == 4 else (WHITE if ci == 0 else LIGHT)
        bld = ci == 4
        add_text(slide, x, y + Inches(0.1), w, row_h,
                 cell, size=13, color=clr, bold=bld, align=PP_ALIGN.CENTER)

# bottom note
add_text(slide, Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.0),
         '首个实现 "零门槛自然语言交互 + 长效记忆自主学习 + 事件驱动群组协作" 的国产框架\n'
         '支持 MCP 插件扩展接入真实业务场景',
         size=15, color=GREEN, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 8 - User Scenarios
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_line(slide, Inches(0.8), Inches(0.6), Inches(1.5), ACCENT)
add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.6),
         "谁会用？怎么用？", size=36, color=WHITE, bold=True)

scenarios = [
    ("个人开发者", ACCENT, [
        '"帮我写一个快排"',
        "管家直接完成，零配置",
        "开箱即用",
    ]),
    ("小型团队", GREEN, [
        '"创建一个前端专家"',
        "定制专属 Agent",
        "越用越懂你的工作习惯",
    ]),
    ("复杂项目", ORANGE, [
        '"组建开发团队做Todo应用"',
        "需求→架构→编码→测试→文档",
        "全流程 AI 协作",
    ]),
]
for i, (title, clr, items) in enumerate(scenarios):
    x = Inches(0.8 + i * 4.1)
    add_rect(slide, x, Inches(1.6), Inches(3.8), Inches(5.0), BG_CARD, clr)
    add_text(slide, x + Inches(0.3), Inches(1.9), Inches(3.2), Inches(0.5),
             title, size=22, color=clr, bold=True, align=PP_ALIGN.CENTER)
    add_line(slide, x + Inches(0.6), Inches(2.5), Inches(2.6), clr, Pt(2))
    for j, item in enumerate(items):
        add_text(slide, x + Inches(0.4), Inches(2.9 + j * 1.0), Inches(3.0), Inches(0.8),
                 item, size=15, color=LIGHT, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 9 - Business Model
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_line(slide, Inches(0.8), Inches(0.6), Inches(1.5), ACCENT)
add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.6),
         "商业模式", size=36, color=WHITE, bold=True)

phases = [
    ("阶段一", "开源社区版", ACCENT, [
        "免费使用，积累用户和口碑",
        "GitHub 开源，社区驱动迭代",
        "建立开发者生态基础",
    ]),
    ("阶段二", "企业版订阅", GREEN, [
        "行业定制 Agent 模板",
        "金融 / 医疗 / 教育",
        "企业级安全合规 · 审计日志 · SSO",
    ]),
    ("阶段三", "Agent 市场", ORANGE, [
        "用户创建的 Agent/Skill 上架交易",
        "平台抽佣，形成生态",
        "开发者与用户双向价值",
    ]),
]
for i, (phase, title, clr, items) in enumerate(phases):
    x = Inches(0.8 + i * 4.1)
    add_rect(slide, x, Inches(1.6), Inches(3.8), Inches(5.0), BG_CARD, clr)
    add_text(slide, x + Inches(0.3), Inches(1.8), Inches(3.2), Inches(0.35),
             phase, size=13, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.3), Inches(2.2), Inches(3.2), Inches(0.5),
             title, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_line(slide, x + Inches(0.6), Inches(2.8), Inches(2.6), clr, Pt(2))
    for j, item in enumerate(items):
        add_text(slide, x + Inches(0.4), Inches(3.2 + j * 0.9), Inches(3.0), Inches(0.7),
                 f"•  {item}", size=14, color=LIGHT)

# ============================================================
# Slide 10 - Market Analysis
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_line(slide, Inches(0.8), Inches(0.6), Inches(1.5), ACCENT)
add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.6),
         "市场分析", size=36, color=WHITE, bold=True)

# Market size - left
add_rect(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.0), BG_CARD)
add_text(slide, Inches(1.2), Inches(1.8), Inches(4.8), Inches(0.4),
         "市场规模", size=20, color=GREEN, bold=True)

market_data = [
    ("2025年全球AI Agent市场", "65亿美元"),
    ("年复合增长率", "42%"),
    ("2028年预计规模", "260亿美元"),
    ("中国市场占比", "~25%"),
]
for i, (label, val) in enumerate(market_data):
    y = Inches(2.5 + i * 0.9)
    add_rect(slide, Inches(1.2), y, Inches(4.8), Inches(0.7), RGBColor(0x0A, 0x25, 0x40))
    add_text(slide, Inches(1.5), y + Inches(0.1), Inches(2.8), Inches(0.5),
             label, size=14, color=LIGHT)
    add_text(slide, Inches(4.2), y + Inches(0.1), Inches(1.5), Inches(0.5),
             val, size=18, color=ACCENT, bold=True, align=PP_ALIGN.RIGHT)

# Competitors - right
add_rect(slide, Inches(6.8), Inches(1.6), Inches(5.5), Inches(5.0), BG_CARD)
add_text(slide, Inches(7.2), Inches(1.8), Inches(4.8), Inches(0.4),
         "竞品格局", size=20, color=ORANGE, bold=True)

competitors = [
    ("AutoGPT", "GitHub 170K star\n功能单一，无协作能力", RED),
    ("CrewAI", "融资 $18M\n绑定 OpenAI，预设流程", RED),
    ("国内", "暂无成熟的多Agent协作框架", GRAY),
    ("CoBeing", "国产化 + 多厂商 + 完整GUI\n填补市场空白", GREEN),
]
for i, (name, desc, clr) in enumerate(competitors):
    y = Inches(2.5 + i * 1.0)
    add_rect(slide, Inches(7.2), y, Inches(4.8), Inches(0.8), RGBColor(0x0A, 0x25, 0x40))
    add_text(slide, Inches(7.5), y + Inches(0.05), Inches(1.5), Inches(0.35),
             name, size=14, color=clr, bold=True)
    add_text(slide, Inches(7.5), y + Inches(0.35), Inches(4.2), Inches(0.4),
             desc, size=11, color=LIGHT)

# ============================================================
# Slide 11 - Team
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_line(slide, Inches(0.8), Inches(0.6), Inches(1.5), ACCENT)
add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.6),
         "关于我们", size=36, color=WHITE, bold=True)

add_rect(slide, Inches(2.5), Inches(1.8), Inches(8), Inches(5.0), BG_CARD, ACCENT)

items = [
    ("团队规模", "个人开发者"),
    ("项目投入", "1-2个月密集开发"),
    ("技术栈", "TypeScript 全栈，前后端 + AI 工程能力"),
    ("代码规模", "~25K 行代码，417 个测试用例"),
    ("为什么能做成", "深入调研竞品，从架构设计到代码实现全栈把控"),
]
for i, (label, val) in enumerate(items):
    y = Inches(2.2 + i * 0.9)
    add_circle(slide, Inches(3.0), y + Inches(0.05), Inches(0.35), ACCENT)
    add_text(slide, Inches(3.6), y, Inches(2.0), Inches(0.4),
             label, size=16, color=ACCENT, bold=True)
    add_text(slide, Inches(5.6), y, Inches(4.5), Inches(0.4),
             val, size=16, color=LIGHT)

# ============================================================
# Slide 12 - Progress
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_line(slide, Inches(0.8), Inches(0.6), Inches(1.5), ACCENT)
add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.6),
         "项目进展", size=36, color=WHITE, bold=True)

add_rect(slide, Inches(0.8), Inches(1.5), Inches(7.5), Inches(5.2), BG_CARD)
add_text(slide, Inches(1.2), Inches(1.7), Inches(6.8), Inches(0.4),
         "已完成", size=20, color=GREEN, bold=True)

completed = [
    "核心框架开发完成（~25K行代码）",
    "7家国产大模型接入（DeepSeek/智谱/通义/MiniMax/豆包/Moonshot/MiMo）",
    "完整桌面GUI（React 19 + Tauri 2.0）",
    "群组协作全流程可用（任务分解/依赖链/投票/经验沉淀）",
    "Docker 安全沙箱 + 5级权限系统",
    "417 个测试用例全部通过",
    "QQ Bot 渠道对接 + MCP 插件系统",
]
for i, item in enumerate(completed):
    y = Inches(2.3 + i * 0.58)
    add_text(slide, Inches(1.5), y, Inches(6.5), Inches(0.4),
             f"✓  {item}", size=14, color=LIGHT)

# Status box
add_rect(slide, Inches(8.8), Inches(1.5), Inches(3.8), Inches(5.2), RGBColor(0x0A, 0x2E, 0x50), GREEN)
add_text(slide, Inches(9.0), Inches(2.5), Inches(3.4), Inches(0.5),
         "当前状态", size=18, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(9.0), Inches(3.2), Inches(3.4), Inches(0.6),
         "可演示", size=32, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(9.0), Inches(4.0), Inches(3.4), Inches(1.0),
         "核心功能完整可用\n随时可以跑起来", size=14, color=LIGHT, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 13 - Roadmap
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_line(slide, Inches(0.8), Inches(0.6), Inches(1.5), ACCENT)
add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.6),
         "发展规划", size=36, color=WHITE, bold=True)

roadmap = [
    ("短期 · 3个月", ACCENT, [
        "子智能体并行编排（DAG任务图）",
        "更多通信渠道（微信、飞书）",
        "社区版开源发布",
    ]),
    ("中期 · 1年", GREEN, [
        "Agent 市场生态",
        "企业版功能（审计/SSO/权限）",
        "垂直行业 Agent 模板",
    ]),
    ("长期愿景", ORANGE, [
        "国产多Agent协作基础设施",
        "构建 Agent 开发者生态",
        "推动 AI 团队协作普及",
    ]),
]
# timeline line
add_line(slide, Inches(0.8), Inches(3.8), Inches(11.5), RGBColor(0x2A, 0x3D, 0x5A), Pt(4))

for i, (title, clr, items) in enumerate(roadmap):
    x = Inches(0.8 + i * 4.1)
    # dot on timeline
    add_circle(slide, x + Inches(1.5), Inches(3.6), Inches(0.4), clr)
    # card below
    add_rect(slide, x, Inches(4.3), Inches(3.8), Inches(2.8), BG_CARD, clr)
    add_text(slide, x + Inches(0.3), Inches(4.5), Inches(3.2), Inches(0.4),
             title, size=18, color=clr, bold=True, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        add_text(slide, x + Inches(0.4), Inches(5.2 + j * 0.6), Inches(3.0), Inches(0.5),
                 f"•  {item}", size=13, color=LIGHT)

# ============================================================
# Slide 14 - Funding
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_line(slide, Inches(0.8), Inches(0.6), Inches(1.5), ACCENT)
add_text(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.6),
         "资源需求", size=36, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
         "如获奖 / 获得资助，资金将用于：", size=18, color=LIGHT)

funds = [
    ("30%", "云服务器", "演示环境部署与运维", ACCENT),
    ("30%", "LLM API 额度", "测试与模型调优", GREEN),
    ("20%", "社区运营", "推广与开发者运营", ORANGE),
    ("20%", "安全审计", "合规与安全加固", ACCENT2),
]
for i, (pct, title, desc, clr) in enumerate(funds):
    x = Inches(0.8 + i * 3.1)
    add_rect(slide, x, Inches(2.5), Inches(2.8), Inches(3.8), BG_CARD, clr)
    add_text(slide, x + Inches(0.2), Inches(2.8), Inches(2.4), Inches(0.8),
             pct, size=40, color=clr, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(3.7), Inches(2.4), Inches(0.5),
             title, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(4.3), Inches(2.4), Inches(1.0),
             desc, size=13, color=LIGHT, align=PP_ALIGN.CENTER)

# ============================================================
# Slide 15 - Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

# decorative circles
add_circle(slide, Inches(0.5), Inches(0.5), Inches(4), RGBColor(0x0A, 0x2E, 0x50))
add_circle(slide, Inches(9.5), Inches(4.5), Inches(4), RGBColor(0x0D, 0x38, 0x60))

add_line(slide, Inches(4.5), Inches(2.0), Inches(4.3), ACCENT, Pt(4))

add_text(slide, Inches(1), Inches(2.3), Inches(11.3), Inches(1.0),
         "CoBeing", size=56, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(3.3), Inches(11.3), Inches(0.5),
         "让 AI 组队干活", size=28, color=ACCENT, align=PP_ALIGN.CENTER)

# quote
add_rect(slide, Inches(2), Inches(4.2), Inches(9.3), Inches(1.2), BG_CARD)
add_text(slide, Inches(2.3), Inches(4.35), Inches(8.7), Inches(0.9),
         '"未来的AI应用不是更聪明的聊天机器人，而是更默契的AI团队。\n CoBeing正在构建这个未来。"',
         size=16, color=LIGHT, align=PP_ALIGN.CENTER)

# contact
add_text(slide, Inches(1), Inches(5.8), Inches(11.3), Inches(0.4),
         "ch3sh_lc@sjtu.edu.cn", size=16, color=GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(6.3), Inches(11.3), Inches(0.4),
         "谢谢各位评委", size=20, color=WHITE, align=PP_ALIGN.CENTER)

# ============================================================
# Save
# ============================================================
output = r"D:\agent-codes\roadshow\玄武区创新大赛\PPT\模数杯路演.pptx"
import os
os.makedirs(os.path.dirname(output), exist_ok=True)
prs.save(output)
print(f"PPT saved to {output}")
print(f"Total slides: {len(prs.slides)}")
